#!/usr/bin/env python3
"""Convert Marp markdown to editable PPTX with real text elements.

Usage:
    python3 marp_to_pptx.py workshop_slides.md output.pptx

Requires: pip install python-pptx Pillow
"""

import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def parse_marp_markdown(filepath):
    """Parse Marp markdown into slides."""
    with open(filepath) as f:
        text = f.read()

    # Remove YAML frontmatter
    text = re.sub(r'^---\n.*?---\n', '', text, count=1, flags=re.DOTALL)

    # Split on horizontal rules
    raw_slides = re.split(r'\n---\n', text)

    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        # Extract speaker notes from HTML comments
        notes = []
        for m in re.finditer(r'<!--(.*?)-->', raw, flags=re.DOTALL):
            notes.append(m.group(1).strip())
        # Remove HTML comments from content
        clean = re.sub(r'<!--.*?-->', '', raw, flags=re.DOTALL).strip()
        if not clean:
            continue
        slide = parse_slide(clean)
        slide['notes'] = '\n\n'.join(notes)
        slides.append(slide)
    return slides


def parse_slide(raw):
    """Parse a single slide's markdown into structured elements."""
    slide = {'elements': [], 'bg_image': None}
    lines = raw.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i]

        # Background image directive
        bg_match = re.match(r'!\[bg[^\]]*\]\(([^)]+)\)', line)
        if bg_match:
            slide['bg_image'] = bg_match.group(1)
            i += 1
            continue

        # Regular image
        img_match = re.match(r'!\[([^\]]*)\]\(([^)]+)\)', line)
        if img_match:
            slide['elements'].append({
                'type': 'image',
                'alt': img_match.group(1),
                'src': img_match.group(2)
            })
            i += 1
            continue

        # H1 heading
        if line.startswith('# ') and not line.startswith('## '):
            slide['elements'].append({'type': 'h1', 'text': line[2:]})
            i += 1
            continue

        # H2 heading
        if line.startswith('## '):
            slide['elements'].append({'type': 'h2', 'text': line[3:]})
            i += 1
            continue

        # H3 heading
        if line.startswith('### '):
            slide['elements'].append({'type': 'h3', 'text': line[4:]})
            i += 1
            continue

        # Blockquote
        if line.startswith('> '):
            quote_lines = []
            while i < len(lines) and lines[i].startswith('> '):
                quote_lines.append(lines[i][2:])
                i += 1
            slide['elements'].append({'type': 'quote', 'text': '\n'.join(quote_lines)})
            continue

        # Code block
        if line.startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].startswith('```'):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing ```
            slide['elements'].append({'type': 'code', 'text': '\n'.join(code_lines)})
            continue

        # Table
        if '|' in line and i + 1 < len(lines) and re.match(r'\|[-\s|]+\|', lines[i + 1]):
            table_lines = []
            while i < len(lines) and '|' in lines[i]:
                table_lines.append(lines[i])
                i += 1
            slide['elements'].append({'type': 'table', 'lines': table_lines})
            continue

        # Numbered list
        if re.match(r'\d+\.', line):
            items = []
            while i < len(lines) and re.match(r'\d+\.', lines[i]):
                items.append(re.sub(r'^\d+\.\s*', '', lines[i]))
                i += 1
            slide['elements'].append({'type': 'ordered_list', 'items': items})
            continue

        # Bullet list
        if line.startswith('- ') or line.startswith('* '):
            items = []
            while i < len(lines) and (lines[i].startswith('- ') or lines[i].startswith('* ') or lines[i].startswith('  ')):
                if lines[i].startswith('  '):
                    # continuation of previous item
                    if items:
                        items[-1] += '\n' + lines[i].strip()
                else:
                    items.append(lines[i][2:])
                i += 1
            slide['elements'].append({'type': 'bullet_list', 'items': items})
            continue

        # Plain text (non-empty)
        if line.strip():
            slide['elements'].append({'type': 'text', 'text': line.strip()})

        i += 1

    return slide


def apply_inline_formatting(paragraph, text):
    """Apply bold/italic markdown formatting to a paragraph."""
    # Split on bold and italic markers
    parts = re.split(r'(\*\*\*.*?\*\*\*|\*\*.*?\*\*|\*.*?\*|_.*?_)', text)
    for part in parts:
        if part.startswith('***') and part.endswith('***'):
            run = paragraph.add_run()
            run.text = part[3:-3]
            run.font.bold = True
            run.font.italic = True
        elif part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.bold = True
        elif (part.startswith('*') and part.endswith('*')) or (part.startswith('_') and part.endswith('_')):
            run = paragraph.add_run()
            run.text = part[1:-1]
            run.font.italic = True
        else:
            run = paragraph.add_run()
            run.text = part
    return paragraph


def strip_markdown(text):
    """Remove markdown formatting for plain text."""
    text = re.sub(r'\*\*\*(.*?)\*\*\*', r'\1', text)
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'_(.*?)_', r'\1', text)
    text = re.sub(r'`(.*?)`', r'\1', text)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    return text


def is_title_slide(slide_data):
    """Check if this is a title/section slide (only h1, maybe subtitle)."""
    types = [e['type'] for e in slide_data['elements']]
    return 'h1' in types and len([t for t in types if t not in ('h1', 'text')]) == 0


def build_pptx(slides_data, base_dir, output_path):
    """Build a PPTX from parsed slide data."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in slides_data:
        if is_title_slide(slide_data):
            build_title_slide(prs, slide_data, base_dir)
        else:
            build_content_slide(prs, slide_data, base_dir)

    prs.save(output_path)
    print(f"Saved {len(slides_data)} slides to {output_path}")


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


def build_title_slide(prs, slide_data, base_dir):
    """Build a centered title slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    texts = []
    for el in slide_data['elements']:
        if el['type'] == 'h1':
            texts.append(('h1', el['text']))
        elif el['type'] == 'text':
            texts.append(('text', el['text']))

    # Center content vertically
    total_height = sum(Pt(60) if t[0] == 'h1' else Pt(30) for t in texts)
    start_y = (SLIDE_HEIGHT - total_height) // 2

    y = start_y
    for kind, text in texts:
        txBox = slide.shapes.add_textbox(Inches(1), y, SLIDE_WIDTH - Inches(2), Pt(60) if kind == 'h1' else Pt(36))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        apply_inline_formatting(p, text)
        if kind == 'h1':
            for run in p.runs:
                run.font.size = Pt(44)
                run.font.bold = True
            y += Pt(60)
        else:
            for run in p.runs:
                run.font.size = Pt(24)
                run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            y += Pt(36)

    add_speaker_notes(slide, slide_data.get('notes', ''))


def build_content_slide(prs, slide_data, base_dir):
    """Build a content slide with heading and body elements."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    y = Inches(0.5)
    margin_left = Inches(0.75)
    content_width = SLIDE_WIDTH - Inches(1.5)

    # Check for bg image — place it on the right
    bg_image = slide_data.get('bg_image')
    if bg_image:
        img_path = os.path.join(base_dir, bg_image)
        if os.path.exists(img_path):
            img_width = Inches(3.5)
            slide.shapes.add_picture(
                img_path,
                SLIDE_WIDTH - img_width - Inches(0.5),
                Inches(1),
                img_width
            )
            content_width = SLIDE_WIDTH - img_width - Inches(2)

    for el in slide_data['elements']:
        el_type = el['type']

        if el_type == 'h2':
            txBox = slide.shapes.add_textbox(margin_left, y, content_width, Pt(48))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            apply_inline_formatting(p, el['text'])
            for run in p.runs:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            y += Inches(0.75)

        elif el_type == 'h3':
            txBox = slide.shapes.add_textbox(margin_left, y, content_width, Pt(36))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            apply_inline_formatting(p, el['text'])
            for run in p.runs:
                run.font.size = Pt(28)
                run.font.bold = True
            y += Inches(0.55)

        elif el_type == 'text':
            txBox = slide.shapes.add_textbox(margin_left, y, content_width, Pt(28))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            apply_inline_formatting(p, el['text'])
            for run in p.runs:
                run.font.size = Pt(20)
            y += Inches(0.4)

        elif el_type == 'quote':
            txBox = slide.shapes.add_textbox(margin_left + Inches(0.3), y, content_width - Inches(0.3), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, line in enumerate(el['text'].split('\n')):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                apply_inline_formatting(p, line)
                for run in p.runs:
                    run.font.size = Pt(20)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
            y += Inches(0.35) * (len(el['text'].split('\n')) + 1)

        elif el_type in ('bullet_list', 'ordered_list'):
            txBox = slide.shapes.add_textbox(margin_left, y, content_width, Inches(3))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, item in enumerate(el['items']):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                prefix = f"{i+1}. " if el_type == 'ordered_list' else "\u2022 "
                apply_inline_formatting(p, prefix + item)
                for run in p.runs:
                    run.font.size = Pt(20)
                p.space_after = Pt(8)
            y += Inches(0.4) * len(el['items'])

        elif el_type == 'code':
            txBox = slide.shapes.add_textbox(margin_left, y, content_width, Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            # Add background color
            from pptx.oxml.ns import qn
            solidFill = txBox.fill
            solidFill.solid()
            solidFill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            for i, line in enumerate(el['text'].split('\n')):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.size = Pt(14)
                run.font.name = 'Courier New'
            y += Inches(0.25) * max(len(el['text'].split('\n')), 1) + Inches(0.2)

        elif el_type == 'image':
            img_path = os.path.join(base_dir, el['src'])
            if os.path.exists(img_path):
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(img_path)
                    w, h = img.size
                    max_w = content_width
                    max_h = SLIDE_HEIGHT - y - Inches(0.5)
                    scale = min(max_w / Emu(w * 914400 // 96), max_h / Emu(h * 914400 // 96), 1.0)
                    img_w = int(w * 914400 // 96 * scale)
                    img_h = int(h * 914400 // 96 * scale)
                    # Center horizontally
                    x = margin_left + (content_width - img_w) // 2
                    slide.shapes.add_picture(img_path, x, y, img_w, img_h)
                    y += img_h + Inches(0.2)
                except Exception as e:
                    print(f"  Warning: Could not add image {el['src']}: {e}")
            else:
                print(f"  Warning: Image not found: {img_path}")

        elif el_type == 'table':
            table_lines = el['lines']
            # Parse header
            headers = [c.strip() for c in table_lines[0].split('|')[1:-1]]
            # Skip separator line
            rows = []
            for tl in table_lines[2:]:
                cells = [c.strip() for c in tl.split('|')[1:-1]]
                rows.append(cells)

            cols = len(headers)
            if cols > 0:
                col_width = content_width // cols
                tbl_height = Inches(0.35) * (len(rows) + 1)
                table = slide.shapes.add_table(
                    len(rows) + 1, cols,
                    margin_left, y, content_width, tbl_height
                ).table

                for j, h in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = strip_markdown(h)
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(16)
                            run.font.bold = True

                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row):
                        if ci < cols:
                            cell = table.cell(ri + 1, ci)
                            cell.text = strip_markdown(val)
                            for p in cell.text_frame.paragraphs:
                                for run in p.runs:
                                    run.font.size = Pt(14)

                y += tbl_height + Inches(0.2)

    add_speaker_notes(slide, slide_data.get('notes', ''))


if __name__ == '__main__':
    import sys
    src = sys.argv[1]
    dst = sys.argv[2]
    base_dir = os.path.dirname(os.path.abspath(src))
    slides = parse_marp_markdown(src)
    build_pptx(slides, base_dir, dst)
